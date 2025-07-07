from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from PIL import Image
import glob
import os
from pathlib import Path
import numpy as np
import os
import itertools

ROOT_DATASET    = Path(__file__).resolve().parent.parent.parent
results_folder = Path(ROOT_DATASET, 'results')

max_col_num = 1
num_rows = 1

image_top_offset = Inches(0.5)

font_sizes = [24,16]
slide_width = Inches(13.33)
slide_height = Inches(7.5)
margin = Inches(1)
img_slide_pos = np.linspace(0, slide_width, max_col_num + 2)
img_slide_pos = img_slide_pos[1:-1]

# Create a PowerPoint presentation
output_pptx = rf"{results_folder}\group_allresults_20250618.pptx"

prs = Presentation()
prs.slide_height = slide_height
prs.slide_width =  slide_width


fmri_pngs = {
"name"        : "fmri",        
"data_folder" : Path(ROOT_DATASET, 'data', 'fmri', 'derivatives'),
"wild_cards" : [f"*nat*second*glm*.png"],
}
# subject_behav_pngs = {
# "name"        : "behav", 
# "data_folder" : Path(ROOT_DATASET, 'data', 'behav', subject_label, 'single_stream'),
# "run_folders" : True,
# "wild_cards" :  [f"*flag_plot*.png"],
# }
# subject_eyetrack_pngs = {
# "name"        : "eyetrack", 
# "data_folder" : Path(ROOT_DATASET, 'data', 'eyetrack', subject_label, 'single_stream'),
# "run_folders" : True,
# "wild_cards" : [f"recentering*right_eye*.png", f"Trigger.png"],
# }
# subject_motion_correction_pngs = {
# "name"        : "motion_correction", 
# "data_folder" : Path(ROOT_DATASET, 'data', 'fmri', 'derivatives', subject_label, 'motioncorection'),
# "run_folders" : False,
# "wild_cards" : [f"*motion_correction*.png"],
# }

png_types = [fmri_pngs] #subject_behav_pngs, subject_eyetrack_pngs, subject_motion_correction_pngs]

for png_type in png_types:
    print(png_type["name"])
    for wildcard in png_type["wild_cards"]:
        print(wildcard)
        pngs = list(png_type["data_folder"].glob(wildcard))

        for image in pngs:
            print(image.name)       
            slide_layout = prs.slide_layouts[5]  # Blank slide layout
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = image.name
            for p, paragraph in enumerate(slide.shapes.title.text_frame.paragraphs):
                paragraph.font.size = Pt(font_sizes[p])      

            img = Image.open(image)
            img_ratio = img.width / img.height

            # Grid layout cell size
            available_width = slide_width - 2 * margin
            available_height = slide_height - 2 * margin
            cell_width = available_width / max_col_num
            cell_height = available_height / num_rows

            # Resize keeping aspect ratio
            if (cell_width / cell_height) < img_ratio:
                img_width = cell_width
                img_height = img_width / img_ratio
            else:
                img_height = cell_height
                img_width = img_height * img_ratio

            row = 0
            col = 0

            left = margin + col * cell_width + (cell_width - img_width) / 2
            top = margin + row * cell_height + (cell_height - img_height) / 2 + image_top_offset

            slide.shapes.add_picture(str(image), left, top, width=img_width, height=img_height)
    # Save the PowerPoint file
    prs.save(output_pptx)
    print(f"PowerPoint saved as {output_pptx}")