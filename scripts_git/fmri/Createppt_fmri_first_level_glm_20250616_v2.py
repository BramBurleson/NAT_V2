from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from PIL import Image
import glob
import os
from pathlib import Path
import numpy as np
import os
import itertools

ROOT_DATASET    = Path(__file__).resolve().parent.parent.parent
derivatives_folder = Path(ROOT_DATASET, 'data', 'fmri', 'derivatives')
output_pptx = rf"{derivatives_folder}\single_subject_glm_smooth_6_20250618_NAT.pptx"

font_sizes = [24,16]
# Create a PowerPoint presentation
prs = Presentation()

slide_width = Inches(13.33)
slide_height = Inches(7.5)

margin = Inches(1)
prs.slide_height = slide_height
prs.slide_width =  slide_width


subject_labels = ['sub-01', 'sub-02', 'sub-03', 'sub-04', 'sub-05', 'sub-06', 'sub-07', 'sub-08', 'sub-09'] 
# subject_labels = ['sub-01', 'sub-02', 'sub-03', 'sub-04', 'sub-05', 'sub-07', 'sub-08', 'sub-09'] 
# subject_labels = ['sub-09']
max_col_num = 5
num_rows = len(subject_labels)//max_col_num + 1


#%% Creat pptx with each slide 1 contrast and all subjects
# img_slide_pos = np.linspace(0, slide_width, len(subjects) + 2)

img_slide_pos = np.linspace(0, slide_width, max_col_num + 2)
img_slide_pos = img_slide_pos[1:-1]

contrast_names = [
"Allo_Main"   , # "Allo-Explore",
"Ego_Main"    , # "Ego-Explore" ,
"Color_Main"  , # "Color-Explore",
"Ego-Color"   ,
"Allo-Color"  ,
"Allo-Ego"]

# contrast_names =[
#     # "Color_Flags", 
#     # "Ego_Flags",
#     # "Allo_Flags",
#     # "Color_Faces", 
#     # "Ego_Faces",
#     # "Allo_Faces",
#     "Ego-Color_Flags", 
#     "Allo-Color_Flags", 
#     "Allo-Ego_Flags", 
#     "Ego-Color_Faces", 
#     "Allo-Color_Faces", 
#     "Allo-Ego_Faces"]
# "Instructions"]
for c, contrast_name in enumerate(contrast_names):
    print(contrast_name)
    slide_layout = prs.slide_layouts[5]  # Blank slide layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = contrast_name
    for p, paragraph in enumerate(slide.shapes.title.text_frame.paragraphs):
        paragraph.font.size = Pt(font_sizes[p])

    for s, subject_label in enumerate(subject_labels):
        subject_folder = Path(derivatives_folder, subject_label, 'glm')
        row = s // max_col_num
        col = s % max_col_num

        print(subject_label)

        image = list(subject_folder.glob(f'nat*{subject_label}*{contrast_name}*.png'))[0]
        img = Image.open(image)
        img_ratio = img.width / img.height

        # Calculate cell size based on layout
        available_width = slide_width - 2 * margin
        available_height = slide_height - 2 * margin
        cell_width = available_width / max_col_num
        cell_height = available_height / num_rows

        # Resize keeping aspect ratio
        if (cell_width / cell_height) < img_ratio:
            img_width = cell_width
            img_height = img_width / img_ratio
        else:
            img_height = cell_height
            img_width = img_height * img_ratio

        # Center the image in its grid cell
        left = margin + col * cell_width + (cell_width - img_width) / 2
        top = margin + row * cell_height + (cell_height - img_height) / 2

        slide.shapes.add_picture(str(image), left, top, width=img_width, height=img_height)

# Save the PowerPoint file
prs.save(output_pptx)
print(f"PowerPoint saved as {output_pptx}")
    
# Create a PowerPoint presentation
prs = Presentation()
prs.slide_height = slide_height
prs.slide_width =  slide_width
