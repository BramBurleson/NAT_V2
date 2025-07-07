from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from PIL import Image
import glob
import os
from pathlib import Path
import numpy as np
import os
import itertools

ROOT_DATASET    = Path(__file__).resolve().parent.parent.parent
results_folder = Path(ROOT_DATASET, 'results')
subject_labels = ['sub-01', 'sub-02', 'sub-03', 'sub-04', 'sub-05', 'sub-06', 'sub-07', 'sub-08', 'sub-09'] 

max_col_num = 1
num_rows = 1

image_top_offset = Inches(0.5)

font_sizes = [24,16]
slide_width = Inches(13.33)
slide_height = Inches(7.5)
margin = Inches(1)
img_slide_pos = np.linspace(0, slide_width, max_col_num + 2)
img_slide_pos = img_slide_pos[1:-1]

for s, subject_label in enumerate(subject_labels):
    print(subject_label)

    # Create a PowerPoint presentation
    output_pptx = rf"{results_folder}\{subject_label}_allresults_20250617.pptx"

    prs = Presentation()
    prs.slide_height = slide_height
    prs.slide_width =  slide_width
    

    subject_fmri_pngs = {
    "name"        : "fmri",        
    "data_folder" : Path(ROOT_DATASET, 'data', 'fmri', 'derivatives', subject_label, 'glm'),
    "run_folders" : False,
    "wild_cards" : [f"*design*.png", f"*nat*glm*.png"],
    }
    subject_behav_pngs = {
    "name"        : "behav", 
    "data_folder" : Path(ROOT_DATASET, 'data', 'behav', subject_label, 'single_stream'),
    "run_folders" : True,
    "wild_cards" :  [f"*flag_plot*.png"],
    }
    subject_eyetrack_pngs = {
    "name"        : "eyetrack", 
    "data_folder" : Path(ROOT_DATASET, 'data', 'eyetrack', subject_label, 'single_stream'),
    "run_folders" : True,
    "wild_cards" : [f"recentering*right_eye*.png", f"Trigger.png"],
    }
    subject_motion_correction_pngs = {
    "name"        : "motion_correction", 
    "data_folder" : Path(ROOT_DATASET, 'data', 'fmri', 'derivatives', subject_label, 'motioncorection'),
    "run_folders" : False,
    "wild_cards" : [f"*motion_correction*.png"],
    }

    subject_png_types = [subject_fmri_pngs, subject_behav_pngs, subject_eyetrack_pngs, subject_motion_correction_pngs]

    for subject_png_type in subject_png_types:
        print(subject_png_type["name"])
        for wildcard in subject_png_type["wild_cards"]:
            print(wildcard)
            if subject_png_type["run_folders"]:
                run_folders = list(subject_png_type["data_folder"].glob(f"*run*"))
                pngs = [png for run_folder in run_folders for png in run_folder.glob(wildcard)]
            else:
                pngs = list(subject_png_type["data_folder"].glob(wildcard))

            for image in pngs:
                print(image.name)       
                slide_layout = prs.slide_layouts[5]  # Blank slide layout
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = image.name
                for p, paragraph in enumerate(slide.shapes.title.text_frame.paragraphs):
                    paragraph.font.size = Pt(font_sizes[p])      

                img = Image.open(image)
                img_ratio = img.width / img.height

                # Grid layout cell size
                available_width = slide_width - 2 * margin
                available_height = slide_height - 2 * margin
                cell_width = available_width / max_col_num
                cell_height = available_height / num_rows

                # Resize keeping aspect ratio
                if (cell_width / cell_height) < img_ratio:
                    img_width = cell_width
                    img_height = img_width / img_ratio
                else:
                    img_height = cell_height
                    img_width = img_height * img_ratio

                row = 0
                col = 0

                left = margin + col * cell_width + (cell_width - img_width) / 2
                top = margin + row * cell_height + (cell_height - img_height) / 2 + image_top_offset

                slide.shapes.add_picture(str(image), left, top, width=img_width, height=img_height)
    # Save the PowerPoint file
    prs.save(output_pptx)
    print(f"PowerPoint saved as {output_pptx}")