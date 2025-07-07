from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from PIL import Image
from pathlib import Path
import numpy as np
import os

def scale_image(img_slide_pos, img, slide_height, image_width):
    aspect_ratio = img.height / img.width
    image_height = image_width * aspect_ratio
    left_x = (img_slide_pos - image_width/2)
    top_y = (slide_height - image_height) / 2
    return (image_width, image_height, left_x, top_y)

#Creat pptx with each slide 1 image preprocessing step
slide_width = Inches(15)
slide_height = Inches(15)
font_sizes = [24, 16]
prs = Presentation()
prs.slide_height = slide_height
prs.slide_width =  slide_width
slide_layout = prs.slide_layouts[0]  # Title Slide layout
slide = prs.slides.add_slide(slide_layout)
tmp = list(slide.shapes)
tmp[0].name
tmp[1].name

ROOT_DATASET =  Path(__file__).resolve().parent.parent.parent
print(ROOT_DATASET)
eyetrack_folder = Path(ROOT_DATASET, 'data', 'eyetrack')
output_pptx     = rf"{eyetrack_folder}\gaze_by_block_fractions.pptx"

imgs_per_slide = 1
img_slide_pos = np.linspace(0, slide_width, imgs_per_slide + 2)
img_slide_pos = img_slide_pos[1:-1]
image_width = slide_width / imgs_per_slide + 2

slides = []
subjects = eyetrack_folder.glob("sub*")
for subject in subjects:
    runs = subject.glob("run*")
    for run in runs:
        block_gaze_folder = list(run.glob("*block_gaze_trajectory*"))[0]
        block_gaze_files = list(block_gaze_folder.glob("Gaze_in_block_fractions*"))
        for block_gaze_file in block_gaze_files:  
            file_path = Path(block_gaze_file)
            slide_layout = prs.slide_layouts[5]  # Blank slide layout
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = file_path.name
            # slide.shapes.placeholders[1].text = "threshold +/- 5% > chance slides -- between subjects"
            # slide.placeholders[1].text_frame.paragraphs[0].font.size = Pt(font_sizes[1])  
            print(file_path.name)
            img = Image.open(file_path)
            print(img.width, img.height)
            image_width, image_height, left_x, top_y = scale_image(img_slide_pos, img, slide_height, image_width)
            slide.shapes.add_picture(str(file_path), left_x, top_y, width=image_width, height=image_height)         # Add image to the slide

# # Save the PowerPoint file
prs.save(output_pptx)
print(f"PowerPoint saved as {output_pptx}")