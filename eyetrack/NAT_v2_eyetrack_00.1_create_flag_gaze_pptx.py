from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from PIL import Image
from pathlib import Path
import numpy as np

def scale_image(img_slide_pos, img, slide_height, image_width):
    aspect_ratio = img.height / img.width
    image_height = image_width * aspect_ratio
    left_x = (img_slide_pos - image_width/2)
    top_y = (slide_height - image_height) / 2
    return (image_width, image_height, left_x, top_y)

#Creat pptx with each slide 1 image preprocessing step
slide_width = Inches(15)
slide_height = Inches(15)
font_sizes = [24, 16]
prs = Presentation()
prs.slide_height = slide_height
prs.slide_width =  slide_width
slide_layout = prs.slide_layouts[0]  # Title Slide layout
slide = prs.slides.add_slide(slide_layout)
tmp = list(slide.shapes)
tmp[0].name
tmp[1].name

ROOT_DATASET =  Path(__file__).resolve().parent.parent.parent
print(ROOT_DATASET)
eyetrack_folder = Path(ROOT_DATASET, 'data', 'eyetrack')
output_pptx     = rf"{eyetrack_folder}\flag_gaze_press.pptx"

imgs_per_slide = 1
img_slide_pos = np.linspace(0, slide_width, imgs_per_slide + 2)
img_slide_pos = img_slide_pos[1:-1]
image_width = slide_width / imgs_per_slide + 2

slides = []
subjects = eyetrack_folder.glob("sub*")
for subject in subjects:
    runs = Path(subject, 'single_stream').glob("run*")
    for r, run in enumerate(runs):
        if r == 0:
            file_names = [
            # Path(rf"preproc_steps/recentering_eye_run_data_based_on_central_fixation_{subject.name}_{run.name}.png"),
            # Path(rf"preproc_steps/pupil_metrics_histograms_{subject.name}_{run.name}.png"),
            # Path(rf"preproc_steps/pupil_metrics_timeseries_{subject.name}_{run.name}.png"),
            # Path(rf"preproc_steps/gaze_timeseries_with_blinks_{subject.name}_{run.name}.png"),
            # Path(rf"preproc_steps/gaze_timeseries_blinks_removed_{subject.name}_{run.name}.png"),
            rf"flag_gaze_press_{subject.name}_{run.name}*.png"
            ]
            for file_name in file_names:      
          
                # slide.shapes.title.text = file_name.name
                # slide.shapes.placeholders[1].text = "threshold +/- 5% > chance slides -- between subjects"
                # slide.placeholders[1].text_frame.paragraphs[0].font.size = Pt(font_sizes[1])  

                images = list(Path(run, 'flag_gaze_press_plots').glob(file_name))
                for image in images:
                    slide_layout = prs.slide_layouts[5]  # Blank slide layout
                    slide = prs.slides.add_slide(slide_layout)
                    # print(Path(image).name)
                    img = Image.open(image)
                    # print(img.width, img.height)
                    image_width, image_height, left_x, top_y = scale_image(img_slide_pos, img, slide_height, image_width)
                    slide.shapes.add_picture(str(image), left_x, top_y, width=image_width, height=image_height)         # Add image to the slide

        # # Save the PowerPoint file
    prs.save(output_pptx)
    print(f"PowerPoint saved as {output_pptx}")